"""
Quick processor for document extraction.
Provides a fast preview of document content without extensive processing.
"""

import os
import time
import tempfile
from pathlib import Path

import streamlit as st
from config import config
from processors.base_processor import BaseProcessor, ProcessingResult
from services.azure_openai import answer_with_azure_openai

class QuickProcessor(BaseProcessor):
    """
    Quick document processor that extracts text content quickly
    for a fast preview without extensive processing.
    """
    
    def process_document(self, file_path, doc_dir, images_dir, process_images=True, batch_size=1, status_area=None):
        """
        Quickly extract text from a document for preview.
        
        Args:
            file_path (Path): Path to the document file
            doc_dir (Path): Directory to save processed document
            images_dir (Path): Directory to save extracted images
            process_images (bool): Not used in quick processing
            batch_size (int): Not used in quick processing
            status_area: Streamlit container for status updates
            
        Returns:
            ProcessingResult: Results of the quick extraction
        """
        # Start timer
        start_time = self._start_timer()
        
        # Ensure paths are Path objects
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower()
        
        # Create a temporary directory for output
        temp_dir = Path(tempfile.mkdtemp())
        output_file = temp_dir / f"{file_path.stem}_quick.md"
        
        extracted_text = ""
        try:
            # Extract based on file type
            if file_extension == '.pdf':
                # Use PyPDF2 for quick text extraction from PDF
                try:
                    import PyPDF2
                    with open(file_path, 'rb') as file:
                        reader = PyPDF2.PdfReader(file)
                        num_pages = len(reader.pages)
                        # Add a simple header
                        extracted_text = f"# {file_path.stem}\n\n"
                        # Extract text from each page
                        for page_num in range(num_pages):
                            page = reader.pages[page_num]
                            page_text = page.extract_text()
                            if page_text:
                                extracted_text += f"## Page {page_num + 1}\n\n{page_text}\n\n"
                except ImportError:
                    # Fallback if PyPDF2 is not available
                    extracted_text = "Installing PyPDF2 for faster PDF extraction is recommended.\n\nProcessing with Docling..."
            elif file_extension in ['.docx', '.doc']:
                # Use python-docx for Word documents
                try:
                    import docx
                    doc = docx.Document(file_path)
                    # Add document title
                    extracted_text = f"# {file_path.stem}\n\n"
                    # Extract text from paragraphs
                    for para in doc.paragraphs:
                        if para.text.strip():
                            extracted_text += para.text + "\n\n"
                except ImportError:
                    extracted_text = "Installing python-docx for faster Word document extraction is recommended.\n\nProcessing with Docling..."
            elif file_extension in ['.pptx', '.ppt']:
                # Use python-pptx for PowerPoint
                try:
                    import pptx
                    presentation = pptx.Presentation(file_path)
                    # Add presentation title
                    extracted_text = f"# {file_path.stem}\n\n"
                    # Process slides
                    for i, slide in enumerate(presentation.slides):
                        extracted_text += f"## Slide {i+1}\n\n"
                        # Extract text from shapes
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                extracted_text += shape.text + "\n\n"
                except ImportError:
                    extracted_text = "Installing python-pptx for faster PowerPoint extraction is recommended.\n\nProcessing with Docling..."
            else:
                extracted_text = f"Unsupported file format for quick extraction: {file_extension}\n\nProcessing with Docling..."
        
        except Exception as e:
            extracted_text = f"Error during quick text extraction: {str(e)}\n\nFalling back to Docling processing..."
        
        # Save the extracted text to a file
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(extracted_text)
        
        # Calculate extraction time
        processing_time = self._stop_timer(start_time)
        
        # Return the processing result
        return ProcessingResult(
            content=extracted_text,
            output_file=output_file,
            images_dir=images_dir,  # No images in quick processing
            processing_time=processing_time
        )
    
    def answer_question(self, question, context=None):
        """
        Answer a question based on the quick extract content.
        
        Args:
            question (str): The question to answer
            context (str, optional): Content to use as context (if None, use session state)
            
        Returns:
            str: The answer to the question
        """
        # Use session state content if context not provided
        if context is None and 'quick_extract_content' in st.session_state:
            context = st.session_state.quick_extract_content
        
        if not context:
            return "No document content available to answer questions."
        
        try:
            # Use Azure OpenAI to answer the question
            return answer_with_azure_openai(
                question=question,
                content=context,
                azure_domain=config.AZURE_DOMAIN,
                api_key=config.AZURE_API_KEY,
                deployment_name=config.AZURE_DEPLOYMENT_NAME,
                api_version=config.AZURE_API_VERSION
            )
        except Exception as e:
            return f"Error answering question: {str(e)}"